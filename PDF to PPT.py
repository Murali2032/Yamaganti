import os
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches

# Define the directory containing the PDF files and the output directory
pdf_directory = r"C:\Users\YamagantiMuraliKrish\Music\October 2024"  # Input folder with PDFs
output_directory = r"C:\Users\YamagantiMuraliKrish\Music\PPT output"  # Output folder for PPT files

# Create the output directory if it doesn't exist
if not os.path.exists(output_directory):
    os.makedirs(output_directory)

# Function to process PDFs recursively from all folders and subfolders
def process_pdfs(pdf_directory, output_directory):
    # Walk through the directory and subdirectories
    for root, dirs, files in os.walk(pdf_directory):
        for pdf_file in files:
            if pdf_file.endswith('.pdf'):
                # Full path to the PDF file
                pdf_file_path = os.path.join(root, pdf_file)

                # Create the output PPTX file name using the same name as the PDF (without extension)
                output_ppt = os.path.join(output_directory, pdf_file.replace('.pdf', '.pptx'))

                # Open the PDF file with PyMuPDF
                try:
                    pdf_document = fitz.open(pdf_file_path)
                except Exception as e:
                    print(f"Error opening {pdf_file_path}: {e}")
                    continue  # Skip this file and move on to the next one

                # Create a PowerPoint presentation
                presentation = Presentation()

                # Convert each page to an image and add it to the PowerPoint presentation
                for page_number in range(pdf_document.page_count):
                    # Get the page
                    page = pdf_document.load_page(page_number)

                    # Convert the page to an image (pixmap)
                    pix = page.get_pixmap(dpi=200)  # Higher DPI for better quality

                    # Save the pixmap as an image
                    temp_image = f"temp_page_{page_number + 1}.png"
                    pix.save(temp_image)

                    # Add a slide to the presentation
                    slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # Blank slide layout

                    # Add the image to the slide
                    slide.shapes.add_picture(temp_image, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))

                    # Print progress
                    print(f"Added page {page_number + 1} to the presentation for {pdf_file}.")

                    # Optionally, delete the temporary image to clean up
                    if os.path.exists(temp_image):
                        os.remove(temp_image)

                # Save the PowerPoint presentation with the same name as the PDF
                presentation.save(output_ppt)
                print(f"Presentation for {pdf_file} saved as {output_ppt}.")

# Call the function to start processing
process_pdfs(pdf_directory, output_directory)

print("PDF to PPT conversion completed for all files.")
